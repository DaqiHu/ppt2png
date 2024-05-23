# imports
from collections.abc import Container
from pptx import Presentation
from pptx.util import Inches
import os
import sys

# configs
img_path = r"D:\Path\To\Image\Folder"
output_path = r"D:\Path\To\Release.pptx"

if 1:

    # 处理命令行输入的PPT路径
    if len(sys.argv) == 2:
        args = sys.argv[1:]
        fileName = args[0]

        # 指定图片文件夹路径
        img_path = f"./{fileName}"
        output_path = f"{fileName}.pptx"

    # 创建新的PPT，指定16:9比例
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # 遍历文件夹中的所有图像，每张图片放在单独的幻灯片页面上
    for i, img_file in enumerate(
        sorted(os.listdir(img_path), key=lambda x: int(x.split(".")[0][5:]))
    ):
        # 添加新的幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 添加图片占位符
        pic = slide.shapes.add_picture(
            os.path.join(img_path, img_file),
            Inches(0),
            Inches(0),
            prs.slide_width,
            prs.slide_height,
        )
        # 调整图片大小，保持比例不变，铺满整个幻灯片页面
        pic.width = prs.slide_width
        pic.height = prs.slide_height

    # 保存PPT
    prs.save(output_path)

# debug
if 0:
    import os

    img_path = r"test"
    li = sorted([int(x.split(".")[0][5:]) for x in os.listdir(img_path)])
    print(*li, sep="\n")
